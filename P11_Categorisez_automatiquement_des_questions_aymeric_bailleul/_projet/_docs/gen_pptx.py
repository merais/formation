"""
Script de génération du PowerPoint P11 — RAG Chatbot Puls-Events
Style : gris/anthracite professionnel — version vulgarisée
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette anthracite/gris pro ──────────────────────────────────────────────
C_BG        = RGBColor(0x2B, 0x2B, 0x2B)   # fond slides : anthracite
C_TITLE_BG  = RGBColor(0x1A, 0x1A, 0x1A)   # bande titre : noir doux
C_ACCENT    = RGBColor(0x5B, 0xA3, 0xD0)   # bleu acier (accent)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xD0, 0xD0, 0xD0)   # texte secondaire
C_GREY      = RGBColor(0x8A, 0x8A, 0x8A)   # texte discret
C_GREEN     = RGBColor(0x4C, 0xAF, 0x50)   # ✅ points positifs

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

OUTPUT = os.path.join(os.path.dirname(__file__), "ABAI_P11_presentation.pptx")
PNG_PATH = os.path.join(os.path.dirname(__file__), "pipeline_rag.png")

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]   # complètement vide


def add_slide() -> object:
    return prs.slides.add_slide(blank_layout)


def bg(slide, color: RGBColor = C_BG):
    """Remplit le fond de la slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height, text="", bold=False,
        size=20, color=C_WHITE, bg_color=None, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    """Ajoute une zone de texte."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    return txBox


def hbar(slide, top, color=C_ACCENT, height=0.04):
    """Ligne horizontale décorative."""
    from pptx.util import Inches
    bar = slide.shapes.add_shape(
        1,
        Inches(0), Inches(top), SLIDE_W, Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def title_band(slide, title: str, subtitle: str = ""):
    """Bande titre en haut de chaque slide."""
    rect = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_W, Inches(1.3))
    rect.fill.solid()
    rect.fill.fore_color.rgb = C_TITLE_BG
    rect.line.fill.background()
    box(slide, 0.3, 0.05, 12.5, 0.75, title,
        bold=True, size=26, color=C_ACCENT, align=PP_ALIGN.LEFT)
    if subtitle:
        box(slide, 0.3, 0.78, 12.5, 0.45, subtitle,
            size=14, color=C_LIGHT, align=PP_ALIGN.LEFT)
    hbar(slide, 1.28)


def bullet_lines(slide, left, top, width, height, lines: list,
                 base_size=17):
    """
    lines = [(indent, text, bold, color), ...]
    indent 0 = titre, 1 = bullet, 2 = sous-bullet
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (indent, text, bold, color) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.level = indent
        run = p.add_run()
        run.text = text
        run.font.size = Pt(base_size - indent * 1.5)
        run.font.bold = bold
        run.font.color.rgb = color if color else C_WHITE


def add_table(slide, left, top, width, headers: list, rows: list,
              col_widths=None, header_size=14, body_size=13):
    """Ajoute un tableau simple."""
    from pptx.util import Inches, Pt
    cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(
        n_rows, cols,
        Inches(left), Inches(top), Inches(width), Inches(0.4 * n_rows)).table

    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(width * w / total)

    def fmt_cell(cell, text, hdr=False):
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_TITLE_BG if hdr else RGBColor(0x38, 0x38, 0x38)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(header_size if hdr else body_size)
        run.font.bold = hdr
        run.font.color.rgb = C_ACCENT if hdr else C_WHITE

    for i, h in enumerate(headers):
        fmt_cell(tbl.cell(0, i), h, hdr=True)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            fmt_cell(tbl.cell(r + 1, c), val)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Titre
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
hbar(s, 2.8, C_ACCENT, 0.06)
box(s, 1, 1.2, 11, 1.2,
    "Assistant Conversationnel RAG\npour Puls-Events",
    bold=True, size=36, color=C_WHITE, align=PP_ALIGN.CENTER)
box(s, 1, 3.0, 11, 0.6,
    "Comment interroger un million d\u2019\xe9v\xe9nements culturels en langage naturel\u00a0?",
    size=18, color=C_LIGHT, align=PP_ALIGN.CENTER, italic=True)
box(s, 1, 4.0, 11, 0.5,
    "Aymeric Bailleul  \u00b7  F\xe9vrier 2026  \u00b7  Data Engineer \u2014 OpenClassrooms",
    size=15, color=C_GREY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Contexte (1/2)
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
title_band(s, "Contexte et probl\xe9matique  (1/2)",
           "Le probl\xe8me\u00a0: la recherche classique ne suffit plus")
bullet_lines(s, 0.5, 1.45, 12.3, 5.5, [
    (0, "Puls-Events agr\xe8ge les \xe9v\xe9nements culturels publics de toute la France", False, C_WHITE),
    (1, "Source\u00a0: Open Agenda \u2014 913\u202f818 \xe9v\xe9nements au 03/02/2026", False, C_LIGHT),
    (1, "Concerts, expositions, festivals, spectacles, conf\xe9rences\u2026", False, C_LIGHT),
    (0, "Aujourd\u2019hui\u00a0: une recherche par mots-cl\xe9s retourne des listes brutes", False, C_WHITE),
    (1, "C\u2019est \xe0 l\u2019utilisateur de lire, trier, et comprendre les r\xe9sultats", False, C_LIGHT),
    (0, "Le foss\xe9 entre la question et la r\xe9ponse\u00a0:", True, C_ACCENT),
    (1, "\u00ab\u00a0Y a-t-il des expos gratuites pour enfants \xe0 Montpellier ce weekend\u00a0?\u00a0\u00bb", False, C_LIGHT),
    (1, "\u00ab\u00a0Quel \xe9v\xe9nement en plein air peut-on faire en famille en Occitanie en juillet\u00a0?\u00a0\u00bb", False, C_LIGHT),
    (0, "La recherche par mots-cl\xe9s filtre et liste. Elle ne raisonne pas.", True, C_ACCENT),
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Contexte (2/2)
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
title_band(s, "Contexte et probl\xe9matique  (2/2)",
           "L\u2019objectif\u00a0: un assistant qui comprend et r\xe9pond")
bullet_lines(s, 0.5, 1.45, 12.3, 5.0, [
    (0, "P\xe9rim\xe8tre du POC\u00a0:", True, C_ACCENT),
    (1, "R\xe9gion Occitanie \u2014 p\xe9rim\xe8tre volontairement restreint pour valider la cha\xeene technique", False, C_LIGHT),
    (1, "\xc9v\xe9nements des 12 derniers mois + futurs \u2192 7\u202f960 \xe9v\xe9nements pertinents retenus", False, C_LIGHT),
    (1, "Un POC n\u2019a pas vocation \xe0 \xeatre parfait\u00a0: il prouve que l\u2019approche est viable", False, C_GREY),
    (0, "Question centrale\u00a0:", True, C_ACCENT),
    (1, "Comment permettre \xe0 un utilisateur d\u2019interroger en langage naturel", False, C_WHITE),
    (1, "une base d\u2019\xe9v\xe9nements et obtenir une r\xe9ponse pr\xe9cise, contextualis\xe9e et v\xe9rifiable", False, C_WHITE),
    (1, "\u2014 sans que le syst\xe8me invente des informations\u00a0?", False, C_WHITE),
    (0, "R\xe9ponse\u00a0: l\u2019architecture RAG \u2014 Retrieval-Augmented Generation", True, C_GREEN),
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture RAG : principe
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
title_band(s, "Architecture RAG\u00a0: principe",
           "Le RAG \u2014 chercher d\u2019abord, puis r\xe9diger")

box(s, 0.4, 1.35, 12.3, 0.45,
    "Un LLM classique (ChatGPT\u2026) invente parfois des r\xe9ponses plausibles mais fausses\u00a0: "
    "c\u2019est une hallucination. Le RAG \xe9limine ce risque.",
    size=14, color=C_LIGHT, italic=True)

for x, title, content in [
    (0.4, "\u2460 RETRIEVAL \u2014 Chercher",
     "On transforme la question en\nempreinte num\xe9rique et on cherche\ndans la base les passages\nles plus proches s\xe9mantiquement"),
    (6.8, "\u2461 GENERATION \u2014 R\xe9diger",
     "Le LLM re\xe7oit ces passages\ncomme contexte et r\xe9dige une r\xe9ponse\nen se basant UNIQUEMENT sur eux\n\u2014 pas d\u2019invention possible"),
]:
    rect = s.shapes.add_shape(1, Inches(x), Inches(1.9), Inches(5.8), Inches(3.2))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x38, 0x38, 0x38)
    rect.line.color.rgb = C_ACCENT
    box(s, x+0.15, 2.0, 5.5, 0.55, title, bold=True, size=18, color=C_ACCENT)
    box(s, x+0.15, 2.65, 5.5, 2.4, content, size=15, color=C_WHITE)

bullet_lines(s, 0.4, 5.3, 12.3, 2.0, [
    (0, "Ce que \xe7a garantit\u00a0:", True, C_GREEN),
    (1, "Aucune information invent\xe9e \u2014 le mod\xe8le ne peut r\xe9pondre qu\u2019avec ce qu\u2019on lui fournit", False, C_LIGHT),
    (1, "Les r\xe9ponses sont tra\xe7ables\u00a0: on sait quels passages ont servi", False, C_LIGHT),
    (1, "Le syst\xe8me dit \u2018je ne sais pas\u2019 si l\u2019info n\u2019existe pas dans la base", False, C_LIGHT),
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Pipeline complet (image)
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
title_band(s, "Pipeline complet",
           "De la donn\xe9e brute \xe0 la r\xe9ponse \u2014 vue d\u2019ensemble")

if os.path.exists(PNG_PATH):
    s.shapes.add_picture(PNG_PATH, Inches(0.3), Inches(1.4), Inches(8.5), Inches(5.5))
    left_leg = 9.0
else:
    box(s, 0.5, 2, 8, 5, "[Image pipeline_rag.png non trouv\xe9e]",
        size=14, color=C_GREY)
    left_leg = 4.0

bullet_lines(s, left_leg if os.path.exists(PNG_PATH) else 0.5,
             1.45, 4.1 if os.path.exists(PNG_PATH) else 12.3, 5.5, [
    (0, "Ce que fait chaque \xe9tape\u00a0:", True, C_ACCENT),
    (1, "\u2460 Source\u00a0: fichier Open Agenda brut (913\u202f818 \xe9v\xe9nements)", False, C_LIGHT),
    (1, "\u2461 Preprocessing\u00a0: filtre g\xe9o + temporel \u2192 7\u202f960 utiles", False, C_LIGHT),
    (1, "\u2462 Chunking\u00a0: d\xe9coupage en blocs de texte de taille fixe", False, C_LIGHT),
    (1, "\u2463 Vectorisation\u00a0: chaque bloc \u2192 empreinte 1024D (Mistral)", False, C_LIGHT),
    (1, "\u2464 FAISS\u00a0: stockage de toutes les empreintes (40.94 MB)", False, C_LIGHT),
    (1, "\u2465 RAG\u00a0: recherche + g\xe9n\xe9ration de la r\xe9ponse", False, C_LIGHT),
    (1, "\u2466 Interface Streamlit + \xe9valuation Ragas", False, C_LIGHT),
], base_size=15)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Choix FAISS
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
title_band(s, "Choix techniques\u00a0: FAISS",
           "La biblioth\xe8que de recherche vectorielle")
box(s, 0.4, 1.35, 12.3, 0.55,
    "FAISS (Meta) : quand l\u2019utilisateur pose une question, on la transforme en vecteur "
    "et FAISS trouve en quelques ms les passages s\xe9mantiquement les plus proches dans la base.",
    size=13, color=C_LIGHT, italic=True)
add_table(s, 0.4, 2.0, 12.5,
    ["Crit\xe8re", "FAISS", "Alternatives (Chroma, Weaviate)"],
    [
        ["D\xe9ploiement", "Fichier .bin, z\xe9ro serveur", "Service d\xe9di\xe9 requis"],
        ["Performance", "Recherche exacte < 1ms / 10\u202f480 vecteurs", "Overkill < 100K vecteurs"],
        ["Int\xe9gration LangChain", "FAISS.load_local() natif", "Variable"],
        ["Co\xfbt", "Gratuit, open-source", "Parfois payant en cloud"],
    ],
    col_widths=[2.5, 4, 4])
bullet_lines(s, 0.5, 5.0, 12.3, 1.8, [
    (0, "En chiffres\u00a0: index de 40.94 MB  \u00b7  cr\xe9\xe9 en 0.39 sec  \u00b7  requ\xeate en < 1 milliseconde", False, C_WHITE),
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Choix Mistral AI
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
title_band(s, "Choix techniques\u00a0: Mistral AI",
           "Un seul fournisseur, trois r\xf4les distincts")
box(s, 0.4, 1.35, 12.3, 0.45,
    "Utiliser un seul fournisseur garantit la coh\xe9rence\u00a0: le mod\xe8le qui indexe les textes "
    "est le m\xeame que celui qui comprend les questions \u2014 les vecteurs sont comparables.",
    size=13, color=C_LIGHT, italic=True)
add_table(s, 0.4, 1.9, 12.5,
    ["R\xf4le", "Mod\xe8le", "Ce qu\u2019il fait concr\xe8tement"],
    [
        ["Transformer le texte en vecteur", "mistral-embed",
         "Convertit chaque bloc en 1024 nombres qui encodent son sens"],
        ["G\xe9n\xe9rer la r\xe9ponse", "mistral-small-latest",
         "Lit les passages r\xe9cup\xe9r\xe9s et r\xe9dige une r\xe9ponse en fran\xe7ais"],
        ["\xc9valuer la qualit\xe9", "mistral-large-latest",
         "Joue le r\xf4le de juge pour noter les r\xe9ponses g\xe9n\xe9r\xe9es"],
    ],
    col_widths=[3.5, 3.5, 5.0])
bullet_lines(s, 0.5, 4.6, 12.3, 2.5, [
    (0, "Pattern\u00a0: Producteur l\xe9ger / Juge lourd", True, C_ACCENT),
    (1, "On utilise un mod\xe8le \xe9conomique pour produire des r\xe9ponses en continu", False, C_LIGHT),
    (1, "On r\xe9serve le mod\xe8le premium \u2014 plus co\xfbteux \u2014 uniquement pour l\u2019\xe9valuation ponctuelle", False, C_LIGHT),
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Méthodologie : preprocessing
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
title_band(s, "M\xe9thodologie\u00a0: pr\xe9-processing",
           "De 913\u202f818 \xe0 7\u202f960 \xe9v\xe9nements exploitables")
bullet_lines(s, 0.5, 1.45, 12.3, 5.0, [
    (0, "Pourquoi filtrer\u00a0? Indexer un million d\u2019\xe9v\xe9nements nationaux pour r\xe9pondre sur "
        "l\u2019Occitanie\u00a0= du bruit inutile et des co\xfbts API \xe9lev\xe9s", False, C_GREY),
    (0, "\u2460 Filtre g\xe9ographique", True, C_ACCENT),
    (1, "Ne conserver que les \xe9v\xe9nements Occitanie  \u2192  89\u202f491 \xe9v\xe9nements", False, C_LIGHT),
    (0, "\u2461 Filtre temporel", True, C_ACCENT),
    (1, "\xc9liminer les \xe9v\xe9nements termin\xe9s depuis > 1 an, garder le futur  \u2192  7\u202f960 \xe9v\xe9nements", False, C_LIGHT),
    (0, "\u2462 Nettoyage", True, C_ACCENT),
    (1, "Colonnes avec > 70\u202f% valeurs manquantes supprim\xe9es, balises HTML retir\xe9es", False, C_LIGHT),
    (0, "Champ texte construit par \xe9v\xe9nement (text_for_rag)\u00a0:", True, C_ACCENT),
    (1, "Titre | Description | D\xe9tails | Mots-cl\xe9s | Lieu", False, C_LIGHT),
    (1, "La s\xe9paration | aide le LLM \xe0 distinguer les champs dans ses r\xe9ponses", False, C_GREY),
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Méthodologie : chunking & vectorisation
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
title_band(s, "M\xe9thodologie\u00a0: chunking & vectorisation",
           "D\xe9couper, transformer, indexer")
box(s, 0.4, 1.35, 12.3, 0.45,
    "Un LLM ne peut pas lire un texte trop long d\u2019un coup. "
    "On d\xe9coupe en blocs (chunks) avec chevauchement pour ne pas couper une information \xe0 mi-phrase.",
    size=13, color=C_LIGHT, italic=True)
add_table(s, 0.4, 1.95, 6.5,
    ["Param\xe8tre", "Valeur", "Ce que \xe7a signifie"],
    [
        ["Taille d\u2019un bloc", "250 tokens (~180 mots)",
         "Assez court pour \xeatre pr\xe9cis, assez long pour avoir du sens"],
        ["Chevauchement", "75 tokens (30\u202f%)",
         "Les 75 derniers tokens d\u2019un bloc = les 75 premiers du suivant"],
        ["R\xe9sultat", "10\u202f480 blocs", "En moyenne 1.31 bloc par \xe9v\xe9nement"],
    ],
    col_widths=[2.5, 2.5, 4.5])
bullet_lines(s, 7.2, 1.95, 5.8, 3.0, [
    (0, "Vectorisation\u00a0:", True, C_ACCENT),
    (1, "mistral-embed\u00a0: chaque bloc \u2192 1024 nombres encodant son sens", False, C_LIGHT),
    (1, "Deux textes similaires auront des vecteurs proches", False, C_GREY),
    (1, "Dur\xe9e totale\u00a0: ~4 min  \u00b7  fichier\u00a0: 83 MB", False, C_LIGHT),
])
bullet_lines(s, 0.5, 5.0, 12.3, 2.2, [
    (0, "Retriever MMR \u2014 \xe9viter les r\xe9p\xe9titions\u00a0:", True, C_ACCENT),
    (1, "La recherche simple retournerait plusieurs blocs du m\xeame \xe9v\xe9nement", False, C_LIGHT),
    (1, "MMR s\xe9lectionne des blocs pertinents ET diversifi\xe9s "
        "(70\u202f% pertinence / 30\u202f% nouveaut\xe9)", False, C_LIGHT),
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Résultats Ragas : global
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
title_band(s, "R\xe9sultats Ragas\u00a0: m\xe9triques globales",
           "\xc9valuation objective \u2014 comment mesurer un syst\xe8me RAG\u00a0?")
box(s, 0.4, 1.35, 12.3, 0.45,
    "Ragas utilise un LLM juge (mistral-large-latest) pour noter automatiquement les r\xe9ponses "
    "selon 4 axes. Protocole\u00a0: 5 questions repr\xe9sentatives.",
    size=13, color=C_LIGHT, italic=True)
add_table(s, 0.4, 1.95, 12.5,
    ["M\xe9trique", "Score", "Ce que \xe7a mesure en clair"],
    [
        ["faithfulness", "0.764",
         "Les r\xe9ponses s\u2019appuient-elles vraiment sur les passages r\xe9cup\xe9r\xe9s\u00a0? (pas d\u2019invention)"],
        ["answer_relevancy", "0.910",
         "La r\xe9ponse r\xe9pond-elle bien \xe0 la question pos\xe9e\u00a0?"],
        ["context_precision", "0.700",
         "Les passages r\xe9cup\xe9r\xe9s sont-ils bien tous utiles\u00a0? (peu de bruit)"],
        ["context_recall", "0.583",
         "A-t-on r\xe9cup\xe9r\xe9 toutes les informations n\xe9cessaires pour r\xe9pondre\u00a0?"],
    ],
    col_widths=[3.5, 2, 6])
bullet_lines(s, 0.5, 5.2, 12.3, 2.0, [
    (0, "Point fort\u00a0: answer_relevancy = 0.910 \u2014 le LLM formule des r\xe9ponses de qualit\xe9", False, C_GREEN),
    (0, "Point faible\u00a0: context_recall = 0.583 \u2014 on ne r\xe9cup\xe8re pas toujours tous les passages utiles avec 10 blocs", False, C_WHITE),
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Résultats Ragas : par question
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
title_band(s, "R\xe9sultats Ragas\u00a0: analyse par question", "")
add_table(s, 0.4, 1.4, 12.5,
    ["Question", "faith.", "ans_rel.", "ctx_prec.", "ctx_rec."],
    [
        ["Expositions \xe0 Montpellier\u00a0?", "0.429", "0.895", "1.000", "0.250"],
        ["Spectacles enfants Occitanie\u00a0?", "0.813", "0.911", "1.000", "0.333"],
        ["Festivals musique \xe9t\xe9 Occitanie", "0.857", "0.921", "0.000", "1.000"],
        ["Carcassonne ce weekend\u00a0?", "0.870", "0.922", "1.000", "0.667"],
        ["\xc9v\xe9nements plein air Occitanie", "0.852", "0.903", "0.500", "0.667"],
    ],
    col_widths=[5, 1.5, 1.5, 1.5, 1.5])
bullet_lines(s, 0.5, 5.4, 12.3, 1.8, [
    (0, "Cas int\xe9ressant \u2014 Festivals musique\u00a0: context_precision=0 mais context_recall=1.0", True, C_ACCENT),
    (1, "L\u2019info \xe9tait bien dans la base \u2014 mais le MMR a trop diversifi\xe9 sur cette requ\xeate trop g\xe9n\xe9rale", False, C_LIGHT),
    (0, "Le\xe7on\u00a0: les requ\xeates sans crit\xe8re pr\xe9cis de lieu ou de date sont plus difficiles \xe0 traiter", False, C_WHITE),
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Exemples de réponses
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
title_band(s, "D\xe9mo\u00a0: exemples de r\xe9ponses",
           "Deux cas r\xe9els \u2014 dans le p\xe9rim\xe8tre et hors p\xe9rim\xe8tre")

for y, q_color, q, r in [
    (1.45, C_ACCENT,
     "\u2753 Cas 1 \u2014 \u00ab\u00a0Y a-t-il des expositions \xe0 Montpellier\u00a0?\u00a0\u00bb",
     "Oui, il y a plusieurs expositions \xe0 Montpellier. L\u2019exposition Montpellier, regarder "
     "la ville autrement est accessible le 20 septembre 2025. L\u2019exposition Regards sur l\u2019op\xe9ra "
     "pr\xe9sente une cinquantaine de dessins d\u2019\xe9tudiants depuis le 10 mai 2025."),
    (4.3, RGBColor(0xD0, 0x6E, 0x30),
     "\u2753 Cas 2 \u2014 \u00ab\u00a0Concerts \xe0 Paris\u00a0?\u00a0\u00bb",
     "Je n\u2019ai pas d\u2019informations sur des concerts \xe0 Paris dans ma base de donn\xe9es. "
     "Mon p\xe9rim\xe8tre est limit\xe9 aux \xe9v\xe9nements de la r\xe9gion Occitanie."),
]:
    box(s, 0.5, y, 12.3, 0.45, q, bold=True, size=15, color=q_color)
    rect = s.shapes.add_shape(1, Inches(0.5), Inches(y+0.5), Inches(12.3), Inches(1.3))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x38, 0x38, 0x38)
    rect.line.color.rgb = C_GREY
    box(s, 0.7, y+0.55, 12.0, 1.1, r, size=13, color=C_LIGHT, italic=True)
box(s, 0.5, 4.1, 12.3, 0.4,
    "Cas 1\u00a0: le syst\xe8me cite des \xe9v\xe9nements r\xe9els avec dates pr\xe9cises, tir\xe9s directement de la base.",
    size=13, color=C_GREEN)
box(s, 0.5, 5.8, 12.3, 0.4,
    "Cas 2\u00a0: le syst\xe8me ne cherche pas \xe0 inventer \u2014 il reconna\xeet ses limites et le signale.",
    size=13, color=C_GREY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Interface Streamlit
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
title_band(s, "Interface utilisateur \u2014 Streamlit",
           "Un chat dans le navigateur, sans installation")
bullet_lines(s, 0.5, 1.45, 12.3, 5.0, [
    (0, "Ce qu\u2019on a construit\u00a0:", True, C_ACCENT),
    (1, "Interface de type chat accessible dans le navigateur (localhost:8501)", False, C_LIGHT),
    (1, "Construite avec Streamlit \u2014 framework Python pour interfaces data science", False, C_LIGHT),
    (1, "Aucune installation requise c\xf4t\xe9 utilisateur final", False, C_GREY),
    (0, "Fonctionnalit\xe9s\u00a0:", True, C_ACCENT),
    (1, "Zone de saisie libre en langage naturel", False, C_LIGHT),
    (1, "Historique de la conversation durant la session", False, C_LIGHT),
    (1, "R\xe9ponse en quelques secondes, ancr\xe9e dans les donn\xe9es de la base", False, C_LIGHT),
    (0, "Lancement\u00a0: streamlit run src/rag/app.py", True, C_ACCENT),
    (0, "Questions de d\xe9mo sugg\xe9r\xe9es\u00a0:", True, C_ACCENT),
    (1, "\u2460 \u00ab\u00a0Y a-t-il des spectacles pour enfants en Occitanie\u00a0?\u00a0\u00bb", False, C_LIGHT),
    (1, "\u2461 \u00ab\u00a0Que se passe-t-il \xe0 Carcassonne ce weekend\u00a0?\u00a0\u00bb", False, C_LIGHT),
    (1, "\u2462 \u00ab\u00a0Concerts \xe0 Paris\u00a0?\u00a0\u00bb  \u2192 test hors p\xe9rim\xe8tre", False, C_LIGHT),
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Limitations & améliorations
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
title_band(s, "Limitations et am\xe9liorations",
           "Ce POC a des limites \u2014 et c\u2019est normal")
box(s, 0.4, 1.35, 12.3, 0.45,
    "Un POC sert \xe0 valider une approche avant d\u2019investir. "
    "Les limitations identifi\xe9es sont des orientations pour la suite, pas des \xe9checs.",
    size=13, color=C_LIGHT, italic=True)
add_table(s, 0.4, 1.95, 12.5,
    ["Limitation actuelle", "Impact r\xe9el", "Piste d\u2019am\xe9lioration"],
    [
        ["Donn\xe9es fig\xe9es au 03/02/2026",
         "Les nouveaux \xe9v\xe9nements n\u2019apparaissent pas",
         "Pipeline incr\xe9mental hebdomadaire"],
        ["context_recall = 0.583",
         "Infos manquantes sur requ\xeates larges",
         "Passer k \xe0 15-20, filtres metadata"],
        ["Pas de m\xe9moire entre questions",
         "Chaque question est ind\xe9pendante",
         "Historique LangChain"],
        ["P\xe9rim\xe8tre Occitanie uniquement",
         "Requ\xeates hors r\xe9gion non trait\xe9es",
         "Partitionnement FAISS par r\xe9gion"],
        ["D\xe9pendance API Mistral",
         "Pas de mode hors-ligne",
         "Cache des questions r\xe9currentes"],
    ],
    col_widths=[3.8, 4.0, 4.2])
box(s, 0.5, 6.0, 12.3, 0.5,
    "Long terme\u00a0: base vectorielle g\xe9r\xe9e (Weaviate/Pinecone)  \u00b7  "
    "Fine-tuning embeddings  \u00b7  Interface multimodale",
    size=13, color=C_GREY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Conclusion
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
title_band(s, "Conclusion", "Bilan \u2014 ce que ce POC prouve")
bullet_lines(s, 0.5, 1.45, 12.3, 5.0, [
    (0, "Ce que nous avons construit et valid\xe9\u00a0:", True, C_ACCENT),
    (1, "Une cha\xeene compl\xe8te\u00a0: Open Agenda brut \u2192 indexation vectorielle \u2192 interface chat", False, C_LIGHT),
    (1, "Un syst\xe8me \xe9valu\xe9 quantitativement \u2014 pas seulement une d\xe9mo qui a l\u2019air de marcher", False, C_LIGHT),
    (1, "Une architecture pragmatique\u00a0: aucun serveur suppl\xe9mentaire, reproductible depuis le README", False, C_LIGHT),
    (0, "Les trois r\xe9sultats cl\xe9s\u00a0:", True, C_ACCENT),
    (1, "\u2705  answer_relevancy = 0.910 \u2014 r\xe9ponses bien adapt\xe9es aux questions", False, C_GREEN),
    (1, "\u2705  faithfulness = 0.764 \u2014 le syst\xe8me ne fabrique pas d\u2019informations", False, C_GREEN),
    (1, "\u2705  Stack LangChain + FAISS + Mistral AI\u00a0: d\xe9ployable sans infrastructure lourde", False, C_GREEN),
    (0, "La prochaine \xe9tape naturelle\u00a0:", True, C_ACCENT),
    (1, "Passer d\u2019un snapshot statique \xe0 un pipeline de mises \xe0 jour incr\xe9mentales", False, C_LIGHT),
    (1, "\xc9tendre le p\xe9rim\xe8tre g\xe9ographique progressivement vers le niveau national", False, C_LIGHT),
])
hbar(s, 6.8, C_ACCENT, 0.04)
box(s, 0, 6.9, 13.33, 0.5,
    "Aymeric Bailleul  \u00b7  P11 Data Engineer  \u00b7  OpenClassrooms  \u00b7  F\xe9vrier 2026",
    size=12, color=C_GREY, align=PP_ALIGN.CENTER)

# ── Sauvegarde ────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"\u2705  Fichier g\xe9n\xe9r\xe9 : {OUTPUT}")
