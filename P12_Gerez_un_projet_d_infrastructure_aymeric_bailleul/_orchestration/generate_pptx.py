"""
Génère la présentation PPTX du POC Sport Data Solution.
Usage : python generate_pptx.py
Sortie : _docs/ABAI_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Couleurs ─────────────────────────────────────────────────────────────────
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF3, 0xF4)
MEDIUM_GRAY = RGBColor(0x85, 0x92, 0x9E)
TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)
TABLE_HEADER_BG = RGBColor(0x2E, 0x86, 0xC1)
TABLE_ROW_EVEN = RGBColor(0xEB, 0xF5, 0xFB)
TABLE_ROW_ODD = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ─── Helpers ──────────────────────────────────────────────────────────────────

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_text_box(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    return tf


def add_paragraph(tf, text, font_size=16, color=TEXT_DARK, bold=False,
                  alignment=PP_ALIGN.LEFT, space_after=Pt(6), font_name="Calibri",
                  bullet=False, level=0):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    p.level = level
    return p


def add_table(slide, left, top, width, rows, cols, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.4 * rows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    return table


def style_table_cell(cell, text, font_size=12, bold=False, color=TEXT_DARK,
                     fill_color=None, alignment=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color


def add_header_bar(slide, title_text, subtitle_text=None):
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), DARK_BLUE)
    add_shape(slide, Inches(0), Inches(1.15), SLIDE_W, Inches(0.06), ACCENT_BLUE)
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.6),
                 title_text, font_size=30, color=WHITE, bold=True)
    if subtitle_text:
        add_text_box(slide, Inches(0.6), Inches(0.7), Inches(10), Inches(0.4),
                     subtitle_text, font_size=14, color=MEDIUM_GRAY, bold=False)


def add_footer(slide, slide_num, total=19):
    add_text_box(slide, Inches(0.5), Inches(7.05), Inches(5), Inches(0.3),
                 "Sport Data Solution — POC Avantages Sportifs", font_size=9,
                 color=MEDIUM_GRAY)
    add_text_box(slide, Inches(11), Inches(7.05), Inches(2), Inches(0.3),
                 f"Aymeric Bailleul — {slide_num}/{total}", font_size=9,
                 color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)


def add_kpi_card(slide, left, top, width, height, value, label, color=ACCENT_BLUE):
    add_shape(slide, left, top, width, height, WHITE, border_color=color)
    add_shape(slide, left, top, width, Inches(0.06), color)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.15), width - Inches(0.3),
                 Inches(0.5), value, font_size=28, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.65), width - Inches(0.3),
                 Inches(0.4), label, font_size=11, color=MEDIUM_GRAY,
                 alignment=PP_ALIGN.CENTER)


def add_flow_box(slide, left, top, width, height, text, color=ACCENT_BLUE):
    shape = add_shape(slide, left, top, width, height, color)
    shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].space_before = Pt(2)
    return shape


def add_arrow_down(slide, left, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, Inches(0.3), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MEDIUM_GRAY
    shape.line.fill.background()
    return shape


def add_arrow_right(slide, left, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, Inches(0.35), Inches(0.25))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MEDIUM_GRAY
    shape.line.fill.background()
    return shape


# =============================================================================
# SLIDE 1 — Titre
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

# Bande décorative
add_shape(slide, Inches(0), Inches(3.2), SLIDE_W, Inches(0.08), ACCENT_BLUE)
add_shape(slide, Inches(0), Inches(3.32), SLIDE_W, Inches(0.04), ACCENT_GREEN)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
             "Sport Data Solution", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(0.7),
             "POC Avantages Sportifs", font_size=32, color=ACCENT_BLUE, bold=False,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
             "Pipeline de données de bout en bout", font_size=20, color=MEDIUM_GRAY,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
             "Aymeric Bailleul — Data Engineer", font_size=18, color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.7), Inches(11), Inches(0.4),
             "Formation Data Engineer — OpenClassrooms — Mars 2026", font_size=14,
             color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 2 — Sommaire
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Sommaire")
add_footer(slide, 2)

items_left = [
    ("01", "Contexte et objectifs"),
    ("02", "Données sources"),
    ("03", "Règles métier"),
    ("04", "Architecture technique"),
    ("05", "Infrastructure de données"),
    ("06", "Privacy by Design (RGPD)"),
    ("07", "Pipeline ETL"),
]
items_right = [
    ("08", "Transformations dbt"),
    ("09", "Orchestration Kestra"),
    ("10", "Tests de qualité"),
    ("11", "Notifications Slack"),
    ("12", "Dashboard Power BI"),
    ("13", "Résultats clés"),
    ("14", "Démonstration live & Perspectives"),
]

for i, (num, label) in enumerate(items_left):
    y = Inches(1.6) + Inches(i * 0.65)
    add_shape(slide, Inches(1.0), y, Inches(0.55), Inches(0.45), ACCENT_BLUE)
    add_text_box(slide, Inches(1.0), y + Inches(0.02), Inches(0.55), Inches(0.4),
                 num, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.75), y + Inches(0.02), Inches(4.5), Inches(0.4),
                 label, font_size=16, color=TEXT_DARK)

for i, (num, label) in enumerate(items_right):
    y = Inches(1.6) + Inches(i * 0.65)
    add_shape(slide, Inches(7.2), y, Inches(0.55), Inches(0.45), ACCENT_GREEN)
    add_text_box(slide, Inches(7.2), y + Inches(0.02), Inches(0.55), Inches(0.4),
                 num, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.95), y + Inches(0.02), Inches(4.5), Inches(0.4),
                 label, font_size=16, color=TEXT_DARK)


# =============================================================================
# SLIDE 3 — Contexte et objectifs
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Contexte et objectifs", "Comprendre la mission de Juliette")
add_footer(slide, 3)

# L'entreprise
add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.5), LIGHT_GRAY)
tf = add_rich_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.4), Inches(2.3))
add_paragraph(tf, "L'entreprise", font_size=18, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "Sport Data Solution — startup de monitoring sportif", font_size=14)
add_paragraph(tf, "Co-fondateurs : Alexandre (cycliste) et Juliette (marathonienne)", font_size=14)
add_paragraph(tf, "", font_size=8)
add_paragraph(tf, "La mission", font_size=18, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "Récompenser les salariés sportifs avec des avantages concrets", font_size=14)

# Objectifs du POC
add_shape(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5), LIGHT_GRAY)
tf = add_rich_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(2.3))
add_paragraph(tf, "Objectifs du POC", font_size=18, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "✓  Tester la faisabilité technique de la solution", font_size=14)
add_paragraph(tf, "✓  Déterminer les données à collecter", font_size=14)
add_paragraph(tf, "✓  Calculer l'impact financier pour l'entreprise", font_size=14)

# Deux avantages
add_kpi_card(slide, Inches(0.5), Inches(4.5), Inches(5.8), Inches(2.3),
             "5 % du salaire brut", "Prime sportive\nTrajet domicile–bureau sportif", ACCENT_BLUE)
add_kpi_card(slide, Inches(6.8), Inches(4.5), Inches(6), Inches(2.3),
             "5 jours / an", "Journées bien-être\n≥ 15 activités physiques / an", ACCENT_GREEN)


# =============================================================================
# SLIDE 4 — Données sources
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Données sources", "3 sources de données pour le POC")
add_footer(slide, 4)

# RH
add_shape(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.2), LIGHT_GRAY)
add_shape(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(0.5), ACCENT_BLUE)
add_text_box(slide, Inches(0.7), Inches(1.55), Inches(3.4), Inches(0.4),
             "Données RH", font_size=16, color=WHITE, bold=True)
tf = add_rich_text_box(slide, Inches(0.7), Inches(2.2), Inches(3.4), Inches(4.3))
add_paragraph(tf, "Données+RH.xlsx", font_size=13, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "161 salariés — 11 colonnes", font_size=12, bold=True)
add_paragraph(tf, "• ID salarié", font_size=11)
add_paragraph(tf, "• Nom, Prénom, DDN", font_size=11, color=ACCENT_RED)
add_paragraph(tf, "• BU (département)", font_size=11)
add_paragraph(tf, "• Date d'embauche", font_size=11)
add_paragraph(tf, "• Salaire brut", font_size=11)
add_paragraph(tf, "• Type de contrat, CP", font_size=11)
add_paragraph(tf, "• Adresse domicile", font_size=11)
add_paragraph(tf, "• Moyen de déplacement", font_size=11)
add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "⚠ Données personnelles → RGPD", font_size=11, color=ACCENT_RED, bold=True)

# Sport
add_shape(slide, Inches(4.7), Inches(1.5), Inches(3.8), Inches(5.2), LIGHT_GRAY)
add_shape(slide, Inches(4.7), Inches(1.5), Inches(3.8), Inches(0.5), ACCENT_GREEN)
add_text_box(slide, Inches(4.9), Inches(1.55), Inches(3.4), Inches(0.4),
             "Données sportives", font_size=16, color=WHITE, bold=True)
tf = add_rich_text_box(slide, Inches(4.9), Inches(2.2), Inches(3.4), Inches(4.3))
add_paragraph(tf, "Données+Sportive.xlsx", font_size=13, color=ACCENT_GREEN, bold=True)
add_paragraph(tf, "161 lignes — 2 colonnes", font_size=12, bold=True)
add_paragraph(tf, "• ID salarié", font_size=11)
add_paragraph(tf, "• Pratique d'un sport", font_size=11)
add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "Nettoyage appliqué :", font_size=12, bold=True)
add_paragraph(tf, "• 66 NaN (41 %) → « Non déclaré »", font_size=11)
add_paragraph(tf, "• Typo « Runing » → « Running »", font_size=11)
add_paragraph(tf, "• 15 sports distincts", font_size=11)

# Strava
add_shape(slide, Inches(8.9), Inches(1.5), Inches(3.8), Inches(5.2), LIGHT_GRAY)
add_shape(slide, Inches(8.9), Inches(1.5), Inches(3.8), Inches(0.5), ACCENT_ORANGE)
add_text_box(slide, Inches(9.1), Inches(1.55), Inches(3.4), Inches(0.4),
             "Strava (simulé)", font_size=16, color=WHITE, bold=True)
tf = add_rich_text_box(slide, Inches(9.1), Inches(2.2), Inches(3.4), Inches(4.3))
add_paragraph(tf, "generate_strava.py", font_size=13, color=ACCENT_ORANGE, bold=True)
add_paragraph(tf, "2 256 activités — 12 mois", font_size=12, bold=True)
add_paragraph(tf, "• 95 salariés sportifs", font_size=11)
add_paragraph(tf, "• 15 types de sports", font_size=11)
add_paragraph(tf, "• 5-40 activités / salarié", font_size=11)
add_paragraph(tf, "• Seed reproductible (42)", font_size=11)
add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "Colonnes :", font_size=12, bold=True)
add_paragraph(tf, "• id_salarie, date_debut", font_size=11)
add_paragraph(tf, "• type_sport, distance_m", font_size=11)
add_paragraph(tf, "• duree_s, commentaire", font_size=11)


# =============================================================================
# SLIDE 5 — Règles métier
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Règles métier", "Conditions d'éligibilité aux avantages")
add_footer(slide, 5)

# Prime sportive
add_shape(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(5.2), LIGHT_GRAY)
add_shape(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(0.06), ACCENT_BLUE)
tf = add_rich_text_box(slide, Inches(0.8), Inches(1.7), Inches(5.4), Inches(4.8))
add_paragraph(tf, "💰  Prime sportive — 5 % du salaire brut", font_size=20, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "", font_size=8)
add_paragraph(tf, "Condition : le salarié vient au bureau via un mode de déplacement sportif", font_size=14)
add_paragraph(tf, "Source : déclaratif RH (colonne « Moyen de déplacement »)", font_size=14)
add_paragraph(tf, "", font_size=8)
add_paragraph(tf, "Validation des distances (Google Maps API) :", font_size=14, bold=True)

# Table seuils
table = add_table(slide, Inches(0.8), Inches(3.7), Inches(5.4), 3, 2,
                  col_widths=[Inches(3.2), Inches(2.2)])
style_table_cell(table.cell(0, 0), "Mode de déplacement", bold=True, color=WHITE, fill_color=TABLE_HEADER_BG, alignment=PP_ALIGN.CENTER)
style_table_cell(table.cell(0, 1), "Distance max", bold=True, color=WHITE, fill_color=TABLE_HEADER_BG, alignment=PP_ALIGN.CENTER)
style_table_cell(table.cell(1, 0), "Marche / Running", fill_color=TABLE_ROW_ODD)
style_table_cell(table.cell(1, 1), "≤ 15 km", fill_color=TABLE_ROW_ODD, alignment=PP_ALIGN.CENTER)
style_table_cell(table.cell(2, 0), "Vélo / Trottinette / Autres", fill_color=TABLE_ROW_EVEN)
style_table_cell(table.cell(2, 1), "≤ 25 km", fill_color=TABLE_ROW_EVEN, alignment=PP_ALIGN.CENTER)

tf2 = add_rich_text_box(slide, Inches(0.8), Inches(5.2), Inches(5.4), Inches(1.3))
add_paragraph(tf2, "Adresse entreprise : 1362 Av. des Platanes, 34970 Lattes", font_size=12, color=MEDIUM_GRAY)
add_paragraph(tf2, "Fallback haversine si l'API Google Maps échoue", font_size=12, color=MEDIUM_GRAY)

# Bien-être
add_shape(slide, Inches(6.9), Inches(1.5), Inches(6), Inches(5.2), LIGHT_GRAY)
add_shape(slide, Inches(6.9), Inches(1.5), Inches(6), Inches(0.06), ACCENT_GREEN)
tf = add_rich_text_box(slide, Inches(7.2), Inches(1.7), Inches(5.4), Inches(4.8))
add_paragraph(tf, "🏃  5 journées bien-être / an", font_size=20, color=ACCENT_GREEN, bold=True)
add_paragraph(tf, "", font_size=8)
add_paragraph(tf, "Condition : au minimum 15 activités physiques dans l'année", font_size=14)
add_paragraph(tf, "Source : activités Strava (simulées pour le POC)", font_size=14)
add_paragraph(tf, "", font_size=8)
add_paragraph(tf, "Logique de calcul :", font_size=14, bold=True)
add_paragraph(tf, "• Compter les activités Strava par salarié actif", font_size=13)
add_paragraph(tf, "• Si nb_activites ≥ 15 → éligible → 5 jours", font_size=13)
add_paragraph(tf, "• Sinon → 0 jour", font_size=13)
add_paragraph(tf, "", font_size=8)
add_paragraph(tf, "Les paramètres (seuil, nb jours) sont évolutifs", font_size=12, color=MEDIUM_GRAY)
add_paragraph(tf, "selon les résultats du POC (cf. note de cadrage)", font_size=12, color=MEDIUM_GRAY)


# =============================================================================
# SLIDE 6 — Architecture technique
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Architecture technique", "Vue d'ensemble de la stack")
add_footer(slide, 6)

# Stack technique table
table = add_table(slide, Inches(0.5), Inches(1.5), Inches(6), 9, 2,
                  col_widths=[Inches(2.5), Inches(3.5)])
headers = [("Composant", "Outil")]
rows_data = [
    ("Base de données", "PostgreSQL 16 (Docker)"),
    ("ETL / Ingestion", "Python 3.12 (pandas, psycopg2)"),
    ("Transformations SQL", "dbt-postgres 1.10"),
    ("Orchestration", "Kestra (UI web, cron, logs)"),
    ("Distances", "Google Maps API + haversine"),
    ("Tests", "pytest (51) + dbt tests (37)"),
    ("Notifications", "Slack Webhooks HTTP"),
    ("Visualisation", "Power BI (DirectQuery)"),
]
style_table_cell(table.cell(0, 0), "Composant", bold=True, color=WHITE, fill_color=TABLE_HEADER_BG)
style_table_cell(table.cell(0, 1), "Outil", bold=True, color=WHITE, fill_color=TABLE_HEADER_BG)
for i, (comp, outil) in enumerate(rows_data):
    bg = TABLE_ROW_EVEN if i % 2 == 0 else TABLE_ROW_ODD
    style_table_cell(table.cell(i + 1, 0), comp, bold=True, fill_color=bg)
    style_table_cell(table.cell(i + 1, 1), outil, fill_color=bg)

# Services Docker
add_shape(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5.2), LIGHT_GRAY)
tf = add_rich_text_box(slide, Inches(7.3), Inches(1.6), Inches(5.2), Inches(5.0))
add_paragraph(tf, "5 services Docker", font_size=18, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "postgres (port 5433)", font_size=14, bold=True)
add_paragraph(tf, "  Base sport_data + base interne kestra", font_size=12, color=MEDIUM_GRAY)
add_paragraph(tf, "kestra (port 9000)", font_size=14, bold=True)
add_paragraph(tf, "  Orchestrateur — flows, logs, historique", font_size=12, color=MEDIUM_GRAY)
add_paragraph(tf, "kestra-setup (éphémère)", font_size=14, bold=True)
add_paragraph(tf, "  Import auto du flow + nettoyage tutoriaux", font_size=12, color=MEDIUM_GRAY)
add_paragraph(tf, "dbt-docs (port 4080)", font_size=14, bold=True)
add_paragraph(tf, "  Documentation dbt statique (nginx)", font_size=12, color=MEDIUM_GRAY)
add_paragraph(tf, "dbt-docs-perms (init)", font_size=14, bold=True)
add_paragraph(tf, "  Permissions volume partagé (chmod 777)", font_size=12, color=MEDIUM_GRAY)


# =============================================================================
# SLIDE 7 — Infrastructure Medallion
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Infrastructure de données", "Architecture Medallion — 5 schémas PostgreSQL")
add_footer(slide, 7)

schemas = [
    ("raw", ACCENT_ORANGE, "activites_strava\n(2 256 lignes simulées)"),
    ("staging", ACCENT_BLUE, "employes (8 col. anonymisées)\npratiques_declarees\nactivites_strava (dbt)\ncache_distances (API)"),
    ("gold", ACCENT_GREEN, "eligibilite_prime\neligibilite_bien_etre\nimpact_financier"),
    ("rh_prive", ACCENT_RED, "identites (Nom, Prénom, DDN)\nvue_primes_nominatives"),
    ("meta", MEDIUM_GRAY, "pipeline_state (hash/watermark)\npipeline_runs (historique)"),
]

x_start = Inches(0.4)
box_w = Inches(2.35)
box_h = Inches(3.5)
gap = Inches(0.15)

for i, (name, color, content) in enumerate(schemas):
    x = x_start + i * (box_w + gap)
    y = Inches(1.6)
    add_shape(slide, x, y, box_w, Inches(0.5), color)
    add_text_box(slide, x, y + Inches(0.05), box_w, Inches(0.4),
                 name, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide, x, y + Inches(0.5), box_w, box_h - Inches(0.5), LIGHT_GRAY, border_color=color)
    tf = add_rich_text_box(slide, x + Inches(0.1), y + Inches(0.6), box_w - Inches(0.2), box_h - Inches(0.7))
    for line in content.split("\n"):
        add_paragraph(tf, line, font_size=11, color=TEXT_DARK)

# Arrows between schemas
for i in range(4):
    x = x_start + (i + 1) * (box_w + gap) - gap + Inches(0.02)
    add_arrow_right(slide, x, Inches(2.8))

# Description below
tf = add_rich_text_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5))
add_paragraph(tf, "10 tables + 1 vue — Architecture Medallion adaptée au RGPD", font_size=14, color=TEXT_DARK, bold=True)
add_paragraph(tf, "• XLSX = couche raw logique (sur disque, jamais copiés en base) — seules les données Strava passent par raw en base", font_size=12, color=MEDIUM_GRAY)
add_paragraph(tf, "• staging = 1er niveau DB : anonymisé, nettoyé, incrémental (UPSERT + soft-delete)", font_size=12, color=MEDIUM_GRAY)
add_paragraph(tf, "• gold = métriques métier calculées par dbt (éligibilités, impact financier)", font_size=12, color=MEDIUM_GRAY)


# =============================================================================
# SLIDE 8 — RGPD
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Privacy by Design (RGPD)", "Minimisation des données — Art. 5.1.c & Art. 25 RGPD")
add_footer(slide, 8)

# Flow boxes
y = Inches(1.7)
add_flow_box(slide, Inches(0.5), y, Inches(3), Inches(1.2),
             "Fichier XLSX\n(11 colonnes dont\nNom, Prénom, DDN)", ACCENT_ORANGE)
add_arrow_right(slide, Inches(3.7), y + Inches(0.45))
add_flow_box(slide, Inches(4.3), y, Inches(2.5), Inches(1.2),
             "Python\n(séparation\nen mémoire)", MEDIUM_GRAY)

# Two arrows out from Python
add_arrow_right(slide, Inches(7.0), Inches(1.85))
add_flow_box(slide, Inches(7.6), Inches(1.5), Inches(5.3), Inches(0.7),
             "staging.employes — 8 colonnes anonymisées", ACCENT_BLUE)

add_arrow_right(slide, Inches(7.0), Inches(2.6))
add_flow_box(slide, Inches(7.6), Inches(2.35), Inches(5.3), Inches(0.7),
             "rh_prive.identites — Nom, Prénom, DDN", ACCENT_RED)

# Roles table
table = add_table(slide, Inches(0.5), Inches(3.6), Inches(12.3), 4, 4,
                  col_widths=[Inches(2.5), Inches(3.5), Inches(3.5), Inches(2.8)])
style_table_cell(table.cell(0, 0), "Rôle PostgreSQL", bold=True, color=WHITE, fill_color=TABLE_HEADER_BG)
style_table_cell(table.cell(0, 1), "Accès", bold=True, color=WHITE, fill_color=TABLE_HEADER_BG)
style_table_cell(table.cell(0, 2), "Usage", bold=True, color=WHITE, fill_color=TABLE_HEADER_BG)
style_table_cell(table.cell(0, 3), "rh_prive", bold=True, color=WHITE, fill_color=TABLE_HEADER_BG, alignment=PP_ALIGN.CENTER)

roles_data = [
    ("role_pipeline", "raw + staging + gold (R/W)", "Pipeline ETL", "❌ Non"),
    ("role_analytics", "gold (lecture seule)", "Power BI", "❌ Non"),
    ("role_rh_admin", "hérite pipeline + rh_prive", "Slack nominatif", "✅ Oui"),
]
for i, (role, access, usage, rh) in enumerate(roles_data):
    bg = TABLE_ROW_EVEN if i % 2 == 0 else TABLE_ROW_ODD
    style_table_cell(table.cell(i + 1, 0), role, bold=True, fill_color=bg)
    style_table_cell(table.cell(i + 1, 1), access, fill_color=bg)
    style_table_cell(table.cell(i + 1, 2), usage, fill_color=bg)
    style_table_cell(table.cell(i + 1, 3), rh, fill_color=bg, alignment=PP_ALIGN.CENTER)

tf = add_rich_text_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5))
add_paragraph(tf, "Principe du moindre privilège : seul role_rh_admin accède aux données nominatives", font_size=14, color=TEXT_DARK, bold=True)
add_paragraph(tf, "• Power BI ne voit que gold.* → aucune donnée personnelle dans les dashboards publics", font_size=12, color=MEDIUM_GRAY)
add_paragraph(tf, "• Les notifications Slack utilisent rh_prive.identites via role_rh_admin", font_size=12, color=MEDIUM_GRAY)
add_paragraph(tf, "• Soft-delete (actif=FALSE) : suppression RGPD des identités des employés retirés du fichier", font_size=12, color=MEDIUM_GRAY)


# =============================================================================
# SLIDE 9 — Pipeline ETL
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Pipeline ETL", "6 étapes Python — main.py --reset")
add_footer(slide, 9)

steps = [
    ("1", "init_db.py", "Schémas + Tables\n+ 3 Rôles PostgreSQL", DARK_BLUE),
    ("2", "load_rh.py", "XLSX → staging.employes\n+ rh_prive.identites", ACCENT_BLUE),
    ("3", "load_sport.py", "XLSX → staging\n.pratiques_declarees", ACCENT_BLUE),
    ("4", "generate_strava.py", "Simulation\n→ raw.activites_strava", ACCENT_ORANGE),
    ("5", "staging.py", "Nettoyage\nraw → staging", ACCENT_GREEN),
    ("6", "gold.py", "Distances + Éligibilités\n+ Impact financier", ACCENT_GREEN),
]

for i, (num, name, desc, color) in enumerate(steps):
    x = Inches(0.3) + i * Inches(2.1)
    y = Inches(1.8)
    # Number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), y, Inches(0.6), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    p = shape.text_frame.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Name
    add_text_box(slide, x, y + Inches(0.75), Inches(2), Inches(0.4),
                 name, font_size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Description box
    add_shape(slide, x + Inches(0.1), y + Inches(1.2), Inches(1.8), Inches(1.3), LIGHT_GRAY, border_color=color)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.3), Inches(1.6), Inches(1.1),
                 desc, font_size=10, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)

    # Arrow
    if i < 5:
        add_arrow_right(slide, x + Inches(1.95), y + Inches(1.7))

# Key points
tf = add_rich_text_box(slide, Inches(0.5), Inches(5.0), Inches(12), Inches(2.0))
add_paragraph(tf, "Points clés du pipeline", font_size=16, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "• UPSERT + soft-delete (actif = FALSE) : idempotence, gestion des départs sans perte d'historique", font_size=13)
add_paragraph(tf, "• Hash SHA-256 des fichiers sources : détection intelligente des changements", font_size=13)
add_paragraph(tf, "• Watermark Strava (inserted_at) : traitement incrémental des nouvelles activités", font_size=13)
add_paragraph(tf, "• Cache distances en base (staging.cache_distances) : évite les appels API redondants", font_size=13)


# =============================================================================
# SLIDE 10 — dbt
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Transformations dbt", "4 modèles SQL — staging + gold")
add_footer(slide, 10)

table = add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), 5, 4,
                  col_widths=[Inches(3), Inches(1.3), Inches(4), Inches(4)])
style_table_cell(table.cell(0, 0), "Modèle", bold=True, color=WHITE, fill_color=TABLE_HEADER_BG)
style_table_cell(table.cell(0, 1), "Schéma", bold=True, color=WHITE, fill_color=TABLE_HEADER_BG, alignment=PP_ALIGN.CENTER)
style_table_cell(table.cell(0, 2), "Description", bold=True, color=WHITE, fill_color=TABLE_HEADER_BG)
style_table_cell(table.cell(0, 3), "Filtres appliqués", bold=True, color=WHITE, fill_color=TABLE_HEADER_BG)

models_data = [
    ("stg_activites_strava", "staging", "Nettoyage raw → staging", "distance>0, durée>0, date ≤395j, sport valide, actif=true"),
    ("eligibilite_prime", "gold", "Prime 5 % salaire brut", "employes × distances × seuils, actif=true"),
    ("eligibilite_bien_etre", "gold", "5 jours si ≥15 activités", "COUNT activités par salarié actif"),
    ("impact_financier", "gold", "Agrégation par département", "FULL OUTER JOIN primes × bien-être"),
]
for i, (model, schema, desc, filtres) in enumerate(models_data):
    bg = TABLE_ROW_EVEN if i % 2 == 0 else TABLE_ROW_ODD
    style_table_cell(table.cell(i + 1, 0), model, bold=True, fill_color=bg)
    style_table_cell(table.cell(i + 1, 1), schema, fill_color=bg, alignment=PP_ALIGN.CENTER)
    style_table_cell(table.cell(i + 1, 2), desc, fill_color=bg)
    style_table_cell(table.cell(i + 1, 3), filtres, fill_color=bg)

# Tests + docs
tf = add_rich_text_box(slide, Inches(0.5), Inches(4.2), Inches(6), Inches(2.5))
add_paragraph(tf, "37 tests dbt", font_size=18, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "• not_null sur toutes les clés et colonnes critiques", font_size=13)
add_paragraph(tf, "• unique sur les clés primaires (id_salarie)", font_size=13)
add_paragraph(tf, "• accepted_values sur type_sport (15 sports)", font_size=13)
add_paragraph(tf, "• Couverture des 4 modèles", font_size=13)

tf = add_rich_text_box(slide, Inches(7), Inches(4.2), Inches(5.5), Inches(2.5))
add_paragraph(tf, "Documentation auto", font_size=18, color=ACCENT_GREEN, bold=True)
add_paragraph(tf, "• dbt docs generate → nginx (port 4080)", font_size=13)
add_paragraph(tf, "• Catalog + lineage graph", font_size=13)
add_paragraph(tf, "• Descriptions YAML sur chaque colonne", font_size=13)
add_paragraph(tf, "• Mise à jour auto à chaque exécution du flow", font_size=13)


# =============================================================================
# SLIDE 11 — Orchestration Kestra
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Orchestration Kestra", "Flow quotidien — détection, ingestion, transformation, notification")
add_footer(slide, 11)

# Flow as boxes
col1_x = Inches(0.5)
col2_x = Inches(4.5)
col3_x = Inches(8.5)

y = Inches(1.6)
add_flow_box(slide, col1_x, y, Inches(3.5), Inches(0.6), "⏰ Cron quotidien 06h00", DARK_BLUE)
add_arrow_down(slide, col1_x + Inches(1.6), y + Inches(0.65))

y = Inches(2.5)
add_flow_box(slide, col1_x, y, Inches(3.5), Inches(0.6), "check_changes.py (hash + watermark)", ACCENT_BLUE)

y = Inches(3.4)
add_flow_box(slide, col1_x, y, Inches(1.6), Inches(0.55), "reload_rh.py", ACCENT_BLUE)
add_flow_box(slide, col1_x + Inches(1.8), y, Inches(1.7), Inches(0.55), "reload_sport.py", ACCENT_BLUE)
add_text_box(slide, col1_x, y + Inches(0.6), Inches(3.5), Inches(0.3),
             "↑ Conditionnels (si fichier modifié)", font_size=9, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

y = Inches(4.4)
add_text_box(slide, col1_x, y, Inches(3.5), Inches(0.3),
             "Si should_run = true ↓", font_size=11, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

add_flow_box(slide, col2_x, Inches(1.6), Inches(3.5), Inches(0.55), "recalculate_distances.py", ACCENT_ORANGE)
add_arrow_down(slide, col2_x + Inches(1.6), Inches(2.2))
add_flow_box(slide, col2_x, Inches(2.5), Inches(3.5), Inches(0.55), "dbt run (4 modèles)", ACCENT_GREEN)
add_arrow_down(slide, col2_x + Inches(1.6), Inches(3.1))
add_flow_box(slide, col2_x, Inches(3.4), Inches(3.5), Inches(0.55), "dbt test (37 tests)", ACCENT_GREEN)
add_arrow_down(slide, col2_x + Inches(1.6), Inches(4.0))
add_flow_box(slide, col2_x, Inches(4.3), Inches(3.5), Inches(0.55), "dbt docs generate", ACCENT_GREEN)

add_flow_box(slide, col3_x, Inches(1.6), Inches(4.3), Inches(0.55), "notify_slack_activities (1 msg/activité)", RGBColor(0x8E, 0x44, 0xAD))
add_arrow_down(slide, col3_x + Inches(2.0), Inches(2.2))
add_flow_box(slide, col3_x, Inches(2.5), Inches(4.3), Inches(0.55), "update_strava_watermark", MEDIUM_GRAY)
add_arrow_down(slide, col3_x + Inches(2.0), Inches(3.1))
add_flow_box(slide, col3_x, Inches(3.4), Inches(4.3), Inches(0.55), "meta.pipeline_runs (INSERT)", MEDIUM_GRAY)
add_arrow_down(slide, col3_x + Inches(2.0), Inches(4.0))
add_flow_box(slide, col3_x, Inches(4.3), Inches(4.3), Inches(0.55), "Slack #sport-avantages (résumé)", RGBColor(0x8E, 0x44, 0xAD))

# Error handler
add_shape(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.5), LIGHT_GRAY, border_color=ACCENT_RED)
tf = add_rich_text_box(slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.3))
add_paragraph(tf, "⚠  Error handlers (en cas d'échec)", font_size=14, color=ACCENT_RED, bold=True)
add_paragraph(tf, "• meta.pipeline_runs → INSERT status='failed'     • Slack #sport-avantages → notification d'échec avec ID exécution + tâche en erreur", font_size=12)
add_paragraph(tf, "• Secrets Kestra (base64) : POSTGRES_PASSWORD, GOOGLE_MAPS_API_KEY, SLACK_WEBHOOK, SLACK_WEBHOOK_INFO", font_size=11, color=MEDIUM_GRAY)


# =============================================================================
# SLIDE 12 — Tests de qualité
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Tests de qualité", "51 tests pytest + 37 tests dbt = 88 tests au total")
add_footer(slide, 12)

# pytest table
add_text_box(slide, Inches(0.5), Inches(1.5), Inches(3), Inches(0.4),
             "pytest — 51 tests", font_size=18, color=ACCENT_BLUE, bold=True)

table = add_table(slide, Inches(0.5), Inches(2.0), Inches(7.5), 7, 3,
                  col_widths=[Inches(2.5), Inches(0.8), Inches(4.2)])
style_table_cell(table.cell(0, 0), "Module", bold=True, color=WHITE, fill_color=TABLE_HEADER_BG)
style_table_cell(table.cell(0, 1), "Tests", bold=True, color=WHITE, fill_color=TABLE_HEADER_BG, alignment=PP_ALIGN.CENTER)
style_table_cell(table.cell(0, 2), "Couverture", bold=True, color=WHITE, fill_color=TABLE_HEADER_BG)

pytest_data = [
    ("test_distances.py", "10", "Haversine, seuils, éligibilité, adresse invalide"),
    ("test_simulation.py", "9", "Distances, durées, dates, IDs, sports, volumes"),
    ("test_staging.py", "12", "Doublons, salaires, dates, complétude, pratiques"),
    ("test_gold.py", "9", "Primes 5 %, bien-être ≥15, impact financier"),
    ("test_rh_prive.py", "9", "Identités RGPD, FK, Privacy by Design, droits"),
    ("conftest.py", "2 fix.", "Connexion DB session-scoped + fixture identites"),
]
for i, (mod, nb, cov) in enumerate(pytest_data):
    bg = TABLE_ROW_EVEN if i % 2 == 0 else TABLE_ROW_ODD
    style_table_cell(table.cell(i + 1, 0), mod, bold=True, fill_color=bg)
    style_table_cell(table.cell(i + 1, 1), nb, fill_color=bg, alignment=PP_ALIGN.CENTER)
    style_table_cell(table.cell(i + 1, 2), cov, fill_color=bg)

# dbt tests
tf = add_rich_text_box(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(4))
add_paragraph(tf, "dbt — 37 tests", font_size=18, color=ACCENT_GREEN, bold=True)
add_paragraph(tf, "", font_size=8)
add_paragraph(tf, "not_null", font_size=14, bold=True, color=ACCENT_GREEN)
add_paragraph(tf, "Toutes les colonnes critiques", font_size=12)
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "unique", font_size=14, bold=True, color=ACCENT_GREEN)
add_paragraph(tf, "Clés primaires (id_salarie)", font_size=12)
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "accepted_values", font_size=14, bold=True, color=ACCENT_GREEN)
add_paragraph(tf, "type_sport (15 valeurs)", font_size=12)
add_paragraph(tf, "", font_size=8)
add_paragraph(tf, "Couverture : 4 modèles dbt", font_size=12, color=MEDIUM_GRAY)
add_paragraph(tf, "Exécutés à chaque run du flow", font_size=12, color=MEDIUM_GRAY)

# KPI cards
add_kpi_card(slide, Inches(0.5), Inches(5.5), Inches(3.5), Inches(1.2),
             "51 / 51 PASS", "Tests pytest", ACCENT_BLUE)
add_kpi_card(slide, Inches(4.5), Inches(5.5), Inches(3.5), Inches(1.2),
             "37 / 37 PASS", "Tests dbt", ACCENT_GREEN)
add_kpi_card(slide, Inches(8.5), Inches(5.5), Inches(4.3), Inches(1.2),
             "88 / 88 PASS", "Tests totaux", DARK_BLUE)


# =============================================================================
# SLIDE 13 — Notifications Slack
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Notifications Slack", "Deux canaux — activités personnalisées + compte-rendu pipeline")
add_footer(slide, 13)

# Canal 1
add_shape(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(3.3), LIGHT_GRAY, border_color=RGBColor(0x8E, 0x44, 0xAD))
add_shape(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(0.06), RGBColor(0x8E, 0x44, 0xAD))
tf = add_rich_text_box(slide, Inches(0.8), Inches(1.7), Inches(5.4), Inches(3.0))
add_paragraph(tf, "Canal 1 — Messages activités", font_size=18, color=RGBColor(0x8E, 0x44, 0xAD), bold=True)
add_paragraph(tf, "SLACK_WEBHOOK", font_size=12, color=MEDIUM_GRAY)
add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "1 message personnalisé par nouvelle activité Strava", font_size=13)
add_paragraph(tf, "15 sports × 2-3 templates = ~35 messages distincts", font_size=13)
add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "« Bravo Juliette ! Tu viens de courir 10.8 km", font_size=12, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "  en 46 min ! Quelle énergie ! »", font_size=12, color=ACCENT_BLUE, bold=True)

# Canal 2
add_shape(slide, Inches(6.9), Inches(1.5), Inches(6), Inches(3.3), LIGHT_GRAY, border_color=ACCENT_GREEN)
add_shape(slide, Inches(6.9), Inches(1.5), Inches(6), Inches(0.06), ACCENT_GREEN)
tf = add_rich_text_box(slide, Inches(7.2), Inches(1.7), Inches(5.4), Inches(3.0))
add_paragraph(tf, "Canal 2 — Infos pipeline", font_size=18, color=ACCENT_GREEN, bold=True)
add_paragraph(tf, "SLACK_WEBHOOK_INFO", font_size=12, color=MEDIUM_GRAY)
add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "Compte-rendu d'exécution (succès / échec)", font_size=13)
add_paragraph(tf, "Résumé chiffré : nb activités, RH/Sport mis à jour", font_size=13)
add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "✅ Pipeline Sport Data — exécution réussie", font_size=12, color=ACCENT_GREEN, bold=True)
add_paragraph(tf, "❌ Pipeline Sport Data — échec (+ tâche en erreur)", font_size=12, color=ACCENT_RED, bold=True)

# Tech details
tf = add_rich_text_box(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5))
add_paragraph(tf, "Implémentation technique", font_size=16, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "• notify_slack_activities.py : JOIN rh_prive.identites pour nom/prénom réels (connexion superuser)", font_size=12)
add_paragraph(tf, "• Messages aléatoires (random.choice) pour la variété — envoi séquentiel par urllib.request", font_size=12)
add_paragraph(tf, "• Exécuté AVANT update_strava_watermark (pour détecter les nouvelles activités)", font_size=12)


# =============================================================================
# SLIDE 14 — Dashboard Power BI
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Dashboard Power BI", "Connexion DirectQuery — role_analytics (lecture seule gold + staging)")
add_footer(slide, 14)

# Page 1 - KPIs publics
add_shape(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(5.2), LIGHT_GRAY, border_color=ACCENT_BLUE)
add_shape(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(0.5), ACCENT_BLUE)
add_text_box(slide, Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.4),
             "Page 1 — KPIs publics (interne)", font_size=16, color=WHITE, bold=True)
tf = add_rich_text_box(slide, Inches(0.8), Inches(2.2), Inches(5.4), Inches(4.3))
add_paragraph(tf, "• Camembert sportifs / non-sportifs (59 % / 41 %)", font_size=13)
add_paragraph(tf, "• Top 10 sports déclarés (histogramme horizontal)", font_size=13)
add_paragraph(tf, "• Top 5 sportifs (classement par nb activités)", font_size=13)
add_paragraph(tf, "• Carte géographique des salariés", font_size=13)
add_paragraph(tf, "• Taux de participation (3 derniers mois)", font_size=13)
add_paragraph(tf, "", font_size=8)
add_paragraph(tf, "Source : staging.employes + staging.pratiques_declarees", font_size=11, color=MEDIUM_GRAY)
add_paragraph(tf, "         + staging.activites_strava", font_size=11, color=MEDIUM_GRAY)

# Page 2 - KPIs RH
add_shape(slide, Inches(6.9), Inches(1.5), Inches(6), Inches(5.2), LIGHT_GRAY, border_color=ACCENT_GREEN)
add_shape(slide, Inches(6.9), Inches(1.5), Inches(6), Inches(0.5), ACCENT_GREEN)
add_text_box(slide, Inches(7.1), Inches(1.55), Inches(5.6), Inches(0.4),
             "Page 2 — KPIs RH", font_size=16, color=WHITE, bold=True)
tf = add_rich_text_box(slide, Inches(7.2), Inches(2.2), Inches(5.4), Inches(4.3))
add_paragraph(tf, "• Histogramme type de mobilité", font_size=13)
add_paragraph(tf, "• Table distances + temps de trajet estimé", font_size=13)
add_paragraph(tf, "• Total jours bien-être gagnés", font_size=13)
add_paragraph(tf, "• Nb salariés éligibles bien-être", font_size=13)
add_paragraph(tf, "• Prime moyenne par salarié éligible", font_size=13)
add_paragraph(tf, "• Total primes sport", font_size=13)
add_paragraph(tf, "", font_size=8)
add_paragraph(tf, "Source : gold.eligibilite_prime", font_size=11, color=MEDIUM_GRAY)
add_paragraph(tf, "         + gold.eligibilite_bien_etre", font_size=11, color=MEDIUM_GRAY)
add_paragraph(tf, "         + staging.cache_distances", font_size=11, color=MEDIUM_GRAY)


# =============================================================================
# SLIDE 15 — Résultats clés
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Résultats clés", "Chiffres du POC")
add_footer(slide, 15)

# KPI cards - row 1
add_kpi_card(slide, Inches(0.5), Inches(1.6), Inches(2.9), Inches(1.5),
             "161", "Salariés", ACCENT_BLUE)
add_kpi_card(slide, Inches(3.7), Inches(1.6), Inches(2.9), Inches(1.5),
             "2 256", "Activités Strava simulées", ACCENT_BLUE)
add_kpi_card(slide, Inches(6.9), Inches(1.6), Inches(2.9), Inches(1.5),
             "68", "Éligibles prime sportive", ACCENT_GREEN)
add_kpi_card(slide, Inches(10.1), Inches(1.6), Inches(2.7), Inches(1.5),
             "172 482 €", "Total primes", ACCENT_GREEN)

# KPI cards - row 2
add_kpi_card(slide, Inches(0.5), Inches(3.5), Inches(2.9), Inches(1.5),
             "67", "Éligibles bien-être", ACCENT_ORANGE)
add_kpi_card(slide, Inches(3.7), Inches(3.5), Inches(2.9), Inches(1.5),
             "335 jours", "Total jours bien-être", ACCENT_ORANGE)
add_kpi_card(slide, Inches(6.9), Inches(3.5), Inches(2.9), Inches(1.5),
             "159", "Adresses Google Maps", MEDIUM_GRAY)
add_kpi_card(slide, Inches(10.1), Inches(3.5), Inches(2.7), Inches(1.5),
             "5 BU", "Départements couverts", MEDIUM_GRAY)

# KPI cards - row 3 (tests)
add_kpi_card(slide, Inches(0.5), Inches(5.4), Inches(3.95), Inches(1.3),
             "51 / 51", "Tests pytest PASS", ACCENT_BLUE)
add_kpi_card(slide, Inches(4.75), Inches(5.4), Inches(3.95), Inches(1.3),
             "37 / 37", "Tests dbt PASS", ACCENT_GREEN)
add_kpi_card(slide, Inches(9.0), Inches(5.4), Inches(3.8), Inches(1.3),
             "88 / 88", "Total tests PASS", DARK_BLUE)


# =============================================================================
# SLIDE 16 — Démonstration live
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Démonstration live", "Scénario : injection d'activités → exécution flow → vérification")
add_footer(slide, 16)

demo_steps = [
    ("1", "Vérifier l'état initial", "Requête SQL :\nnb activités, watermark", DARK_BLUE),
    ("2", "Simuler 5 activités", "python simulate_\nnew_activities.py", ACCENT_ORANGE),
    ("3", "Exécuter le flow", "Kestra UI :\nbouton Execute", ACCENT_BLUE),
    ("4", "Observer les logs", "Kestra UI :\nlogs temps réel", ACCENT_BLUE),
    ("5", "Vérifier Slack", "Messages nominatifs\n+ résumé pipeline", RGBColor(0x8E, 0x44, 0xAD)),
    ("6", "Rafraîchir Power BI", "DirectQuery :\ndonnées à jour", ACCENT_GREEN),
]

for i, (num, title, desc, color) in enumerate(demo_steps):
    x = Inches(0.3) + i * Inches(2.1)
    y = Inches(1.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), y, Inches(0.6), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    p = shape.text_frame.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, x, y + Inches(0.75), Inches(2), Inches(0.4),
                 title, font_size=11, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide, x + Inches(0.1), y + Inches(1.2), Inches(1.8), Inches(1.1), LIGHT_GRAY, border_color=color)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.3), Inches(1.6), Inches(0.9),
                 desc, font_size=10, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)
    if i < 5:
        add_arrow_right(slide, x + Inches(1.95), y + Inches(1.6))

# Commands
tf = add_rich_text_box(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2))
add_paragraph(tf, "Commandes de la démonstration", font_size=16, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "# 1. État initial", font_size=11, color=MEDIUM_GRAY, font_name="Consolas")
add_paragraph(tf, "docker exec sport_data_postgres psql -U postgres -d sport_data -c \"SELECT ...\"", font_size=11, font_name="Consolas")
add_paragraph(tf, "# 2. Simulation", font_size=11, color=MEDIUM_GRAY, font_name="Consolas")
add_paragraph(tf, "python scripts/simulate_new_activities.py", font_size=11, font_name="Consolas")
add_paragraph(tf, "# 3. Flow Kestra → http://localhost:9000 → Execute", font_size=11, color=MEDIUM_GRAY, font_name="Consolas")


# =============================================================================
# SLIDE 17 — Choix techniques
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Choix techniques justifiés")
add_footer(slide, 17)

table = add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), 11, 2,
                  col_widths=[Inches(3.5), Inches(8.8)])
style_table_cell(table.cell(0, 0), "Choix", bold=True, color=WHITE, fill_color=TABLE_HEADER_BG)
style_table_cell(table.cell(0, 1), "Justification", bold=True, color=WHITE, fill_color=TABLE_HEADER_BG)

choices_data = [
    ("PostgreSQL via Docker", "Simule un environnement de production réaliste, reproductible"),
    ("3 rôles PostgreSQL", "Principe du moindre privilège, conformité RGPD (rh_prive isolé)"),
    ("Google Maps API + haversine", "Distances routières réelles, fallback robuste si clé absente"),
    ("Slack Webhooks réels", "Exigence note de cadrage : publications automatiques dans le channel"),
    ("pytest (51) + dbt (37)", "Double couverture : Python (logique métier) + SQL (transformations)"),
    ("dbt (staging + gold)", "Transformations versionnées, testées, documentées, lineage graph"),
    ("Kestra", "Orchestration visuelle, logs, historique d'exécutions, cron, monitoring UI"),
    ("Power BI DirectQuery", "Données toujours à jour, pas d'import manuel → KPIs temps réel"),
    ("UPSERT + soft-delete", "Idempotence, gestion des départs sans perte d'historique"),
    ("SHA-256 + watermark", "Détection intelligente des changements, pas de re-traitement inutile"),
]
for i, (choix, justif) in enumerate(choices_data):
    bg = TABLE_ROW_EVEN if i % 2 == 0 else TABLE_ROW_ODD
    style_table_cell(table.cell(i + 1, 0), choix, bold=True, fill_color=bg)
    style_table_cell(table.cell(i + 1, 1), justif, fill_color=bg)


# =============================================================================
# SLIDE 18 — Perspectives
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Perspectives", "Évolutions envisagées après le POC")
add_footer(slide, 18)

perspectives = [
    ("Connexion Strava réelle", "Remplacer la simulation par l'API OAuth2 Strava pour les données live", ACCENT_BLUE),
    ("Scaling infrastructure", "Passage à un cluster PostgreSQL ou migration cloud (AWS RDS, GCP)", ACCENT_BLUE),
    ("Alerting avancé", "Seuils d'alerte automatiques sur les anomalies de déclaration", ACCENT_ORANGE),
    ("Paramétrage dynamique", "Taux de prime et seuils modifiables sans redéploiement", ACCENT_ORANGE),
    ("Intégration CI/CD", "Tests automatisés à chaque push GitHub (GitHub Actions)", ACCENT_GREEN),
    ("Extension dashboard", "Vue manager par département, évolution temporelle des KPIs", ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(perspectives):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + col * Inches(6.4)
    y = Inches(1.6) + row * Inches(1.7)
    add_shape(slide, x, y, Inches(6), Inches(1.4), LIGHT_GRAY, border_color=color)
    add_shape(slide, x, y, Inches(0.08), Inches(1.4), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(5.4), Inches(0.4),
                 title, font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.6), Inches(5.4), Inches(0.6),
                 desc, font_size=12, color=TEXT_DARK)


# =============================================================================
# SLIDE 19 — Merci
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_shape(slide, Inches(0), Inches(3.2), SLIDE_W, Inches(0.08), ACCENT_BLUE)
add_shape(slide, Inches(0), Inches(3.32), SLIDE_W, Inches(0.04), ACCENT_GREEN)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "Merci !", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.5),
             "Sport Data Solution — POC Avantages Sportifs", font_size=22,
             color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
             "Aymeric Bailleul — Data Engineer", font_size=20, color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.4),
             "aymericbailleul@gmail.com", font_size=16, color=MEDIUM_GRAY,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.6),
             "Questions ?", font_size=28, color=ACCENT_GREEN, bold=True,
             alignment=PP_ALIGN.CENTER)


# ─── Sauvegarde ───────────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(__file__), "_docs", "ABAI_presentation.pptx")
prs.save(output_path)
print(f"Presentation sauvegardee : {output_path}")
